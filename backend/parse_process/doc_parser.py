#!/usr/bin/env python3
"""
doc_parser.py — Universal Document Text Extractor
==================================================
Extracts text from a wide range of document types:
  - Text-layer PDFs        → pdfplumber
  - Scanned / photo PDFs   → pdf2image + pytesseract (OCR)
  - DOCX files             → python-docx
  - Plain text             → direct read
  - Images (JPG/PNG/etc.)  → pytesseract (OCR)
  - XLSX spreadsheets      → openpyxl
  - PPTX presentations     → python-pptx

Usage
-----
    python doc_parser.py <file> [options]

Options
-------
    -o, --output <path>     Save extracted text to a file (default: print to stdout)
    --ocr                   Force OCR even if a text layer is detected
    --lang <code>           OCR language code (default: eng). E.g. hin, fra, deu
    --dpi <n>               DPI for PDF→image rasterisation (default: 300)
    --pages <range>         Page range for PDF, e.g. "1-5" or "3" (default: all)
    --encoding <enc>        Encoding for plain-text files (default: utf-8)
    --quiet                 Suppress progress messages
    -h, --help              Show this help and exit

"""

import argparse
import io
import os
import sys
import re
import subprocess
import tempfile
import unicodedata
from pathlib import Path
from typing import Optional



# Lazy imports with helpful error messages
def _import(package: str, pip_name: str = None):
    """Import a package and surface a clear install message on failure."""
    import importlib
    try:
        return importlib.import_module(package)
    except ImportError:
        pip_name = pip_name or package
        sys.exit(
            f"[ERROR] Missing dependency: '{package}'.\n"
            f"  Install it with:  pip install {pip_name}\n"
        )


# Windows stdout UTF-8 fix
def _fix_stdout_encoding():
    """
    On Windows the console defaults to cp1252 (or similar), which can't
    represent many Unicode characters that appear in PDFs (e.g. U+2212 −,
    U+2019 ', U+00B0 °).  Reconfigure stdout/stderr to UTF-8 if possible.
    Falls back gracefully on older Python or non-Windows platforms.
    """
    if sys.platform != "win32":
        return
    try:
        if hasattr(sys.stdout, "reconfigure"):
            sys.stdout.reconfigure(encoding="utf-8", errors="replace")
            sys.stderr.reconfigure(encoding="utf-8", errors="replace")
        else:
            # Python 3.6 fallback
            sys.stdout = io.TextIOWrapper(
                sys.stdout.buffer, encoding="utf-8", errors="replace"
            )
            sys.stderr = io.TextIOWrapper(
                sys.stderr.buffer, encoding="utf-8", errors="replace"
            )
    except (AttributeError, io.UnsupportedOperation):
        pass  # running in a context without a real buffer (e.g. IDLE) — skip


def unicode_safe(text: str, mode: str = "keep") -> str:
    """
    Make text safe for restricted-encoding terminals.

    mode='keep'     — return text unchanged (default; works after _fix_stdout_encoding)
    mode='replace'  — swap un-encodable chars for closest ASCII equivalents
    mode='ascii'    — strip every non-ASCII character
    """
    if mode == "keep":
        return text
    if mode == "ascii":
        return text.encode("ascii", errors="ignore").decode("ascii")
    if mode == "replace":
        # NFKD decomposition converts many fancy chars to ASCII base + combining marks
        normalized = unicodedata.normalize("NFKD", text)
        return normalized.encode("ascii", errors="ignore").decode("ascii")
    return text


# Helpers
def log(msg: str, quiet: bool = False):
    if not quiet:
        print(f"[doc_parser] {msg}", file=sys.stderr)


def parse_page_range(spec: Optional[str], total: int) -> list[int]:
    """Convert a page-range string like '1-5' or '3' to a 0-based index list."""
    if not spec:
        return list(range(total))
    spec = spec.strip()
    if "-" in spec:
        start, end = spec.split("-", 1)
        start = max(1, int(start))
        end = min(total, int(end))
        return list(range(start - 1, end))
    else:
        n = int(spec)
        if 1 <= n <= total:
            return [n - 1]
        return []


def _pdf_has_text_layer(path: str) -> bool:
    """
    Return True if the PDF has an embedded font (i.e. a text layer).
    Uses pdffonts (poppler) when available, otherwise falls back to pypdf.
    """
    # Try pdffonts (most reliable)
    try:
        result = subprocess.run(
            ["pdffonts", path],
            capture_output=True, text=True, timeout=15
        )
        lines = result.stdout.strip().splitlines()
        # pdffonts header is 2 lines; fonts follow — any extra line = text layer
        return len(lines) > 2
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Fallback: try pypdf extraction on first page
    try:
        pypdf = _import("pypdf")
        reader = pypdf.PdfReader(path)
        if reader.pages:
            text = reader.pages[0].extract_text() or ""
            return bool(text.strip())
    except Exception:
        pass

    return False


def clean_text(text: str) -> str:
    """Normalise whitespace without losing paragraph breaks."""
    # Replace multiple blank lines with a single one
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip trailing whitespace per line
    lines = [line.rstrip() for line in text.splitlines()]
    return "\n".join(lines).strip()


# Extractors
def extract_text_pdf(path: str, pages: Optional[str], quiet: bool) -> str:
    """Extract text from a PDF that has an embedded text layer."""
    pdfplumber = _import("pdfplumber")
    chunks = []
    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        indices = parse_page_range(pages, total)
        log(f"Text-layer PDF — {total} page(s), extracting {len(indices)} page(s).", quiet)
        for i in indices:
            page = pdf.pages[i]
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(chunks)


def extract_ocr_pdf(path: str, pages: Optional[str], lang: str,
                    dpi: int, quiet: bool) -> str:
    """
    OCR a scanned / photographed PDF.
    Rasterises each page with pdf2image, then runs pytesseract.
    """
    pdf2image = _import("pdf2image")
    pytesseract = _import("pytesseract")
    PIL = _import("PIL", "pillow")

    log(f"Scanned PDF — rasterising at {dpi} DPI for OCR (lang={lang}).", quiet)

    # Determine page range before conversion
    # pdf2image uses 1-based first_page / last_page
    from pypdf import PdfReader
    total = len(PdfReader(path).pages)
    indices = parse_page_range(pages, total)  # 0-based

    first_page = (indices[0] + 1) if indices else 1
    last_page = (indices[-1] + 1) if indices else total

    images = pdf2image.convert_from_path(
        path, dpi=dpi, first_page=first_page, last_page=last_page
    )

    chunks = []
    for rel_i, img in enumerate(images):
        abs_page = first_page + rel_i
        log(f"  OCR page {abs_page}/{last_page} …", quiet)
        text = pytesseract.image_to_string(img, lang=lang)
        if text.strip():
            chunks.append(f"--- Page {abs_page} ---\n{text}")

    return "\n\n".join(chunks)


def extract_image(path: str, lang: str, quiet: bool) -> str:
    """OCR a standalone image file."""
    pytesseract = _import("pytesseract")
    PIL_Image = _import("PIL.Image", "pillow")
    from PIL import Image  # noqa: F401 — already imported above

    log(f"Image file — running OCR (lang={lang}).", quiet)
    img = Image.open(path)
    return pytesseract.image_to_string(img, lang=lang)


def extract_docx(path: str, quiet: bool) -> str:
    """Extract text from a .docx file preserving paragraph structure."""
    docx = _import("docx", "python-docx")
    from docx import Document  # noqa: F401

    log("DOCX file — extracting paragraphs and tables.", quiet)
    doc = docx.Document(path)
    parts = []

    # Paragraphs (includes headings)
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n\n".join(parts)


def extract_txt(path: str, encoding: str, quiet: bool) -> str:
    """Read a plain-text file."""
    log(f"Plain-text file (encoding={encoding}).", quiet)
    with open(path, "r", encoding=encoding, errors="replace") as fh:
        return fh.read()


def extract_xlsx(path: str, quiet: bool) -> str:
    """Extract text from an Excel workbook (.xlsx / .xlsm)."""
    openpyxl = _import("openpyxl")
    log("XLSX file — reading all sheets.", quiet)

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).rstrip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def extract_pptx(path: str, quiet: bool) -> str:
    """Extract text from a PowerPoint presentation (.pptx)."""
    pptx = _import("pptx", "python-pptx")
    from pptx import Presentation  # noqa: F401

    log("PPTX file — extracting slide text.", quiet)
    prs = pptx.Presentation(path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)



# Router


IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp", ".gif"}

def parse_document(
    path: str,
    *,
    force_ocr: bool = False,
    lang: str = "eng",
    dpi: int = 300,
    pages: Optional[str] = None,
    encoding: str = "utf-8",
    quiet: bool = False,
) -> str:
    """
    Detect document type and dispatch to the correct extractor.
    Returns the extracted text as a single string.
    """
    p = Path(path)
    if not p.exists():
        sys.exit(f"[ERROR] File not found: {path}")

    ext = p.suffix.lower()

    #  PDF 
    if ext == ".pdf":
        if force_ocr or not _pdf_has_text_layer(path):
            if not force_ocr:
                log("No text layer detected — switching to OCR.", quiet)
            return extract_ocr_pdf(path, pages, lang, dpi, quiet)
        else:
            return extract_text_pdf(path, pages, quiet)

    #  Images 
    if ext in IMAGE_EXTS:
        return extract_image(path, lang, quiet)

    #  Office formats 
    if ext == ".docx":
        return extract_docx(path, quiet)

    if ext in {".xlsx", ".xlsm"}:
        return extract_xlsx(path, quiet)

    if ext == ".pptx":
        return extract_pptx(path, quiet)

    #  Plain text / markdown / code 
    if ext in {".txt", ".md", ".rst", ".csv", ".tsv", ".log", ".json",
               ".xml", ".html", ".htm", ".yaml", ".yml", ""}:
        return extract_txt(path, encoding, quiet)

    #  Unknown — attempt plain-text read with a warning 
    log(f"Unrecognised extension '{ext}' — attempting plain-text read.", quiet)
    try:
        return extract_txt(path, encoding, quiet)
    except Exception as exc:
        sys.exit(f"[ERROR] Could not read '{path}': {exc}")



# CLI


def build_parser() -> argparse.ArgumentParser:
    ap = argparse.ArgumentParser(
        prog="doc_parser.py",
        description="Extract text from PDFs (including scanned), DOCX, images, and more.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    ap.add_argument("file", help="Path to the document to parse.")
    ap.add_argument(
        "-o", "--output",
        metavar="PATH",
        help="Write extracted text to this file (default: stdout).",
    )
    ap.add_argument(
        "--ocr", action="store_true",
        help="Force OCR even when a text layer is present.",
    )
    ap.add_argument(
        "--lang", default="eng", metavar="CODE",
        help="Tesseract language code (default: eng). E.g. hin, fra, deu.",
    )
    ap.add_argument(
        "--dpi", type=int, default=300, metavar="N",
        help="DPI for PDF rasterisation during OCR (default: 300).",
    )
    ap.add_argument(
        "--pages", metavar="RANGE",
        help='Page range for PDFs: e.g. "1-5" or "3". Default: all pages.',
    )
    ap.add_argument(
        "--encoding", default="utf-8", metavar="ENC",
        help="Character encoding for plain-text files (default: utf-8).",
    )
    ap.add_argument(
        "--replace", action="store_true",
        help=(
            "Transliterate non-ASCII characters to their closest ASCII equivalent "
            "before printing (e.g. − → -, ' → '). Useful on restricted terminals."
        ),
    )
    ap.add_argument(
        "--ascii", action="store_true",
        help="Strip all non-ASCII characters from the output. Stronger than --replace.",
    )
    ap.add_argument(
        "--quiet", "-q", action="store_true",
        help="Suppress progress messages.",
    )
    return ap


def main(): # not used with APIs parser_save is used instead
    # Fix Windows cp1252 console before anything is printed
    _fix_stdout_encoding()

    ap = build_parser()
    args = ap.parse_args()

    # Determine Unicode handling mode
    if args.ascii:
        umode = "ascii"
    elif args.replace:
        umode = "replace"
    else:
        umode = "keep"

    text = parse_document(
        args.file,
        force_ocr=args.ocr,
        lang=args.lang,
        dpi=args.dpi,
        pages=args.pages,
        encoding=args.encoding,
        quiet=args.quiet,
    )

    text = clean_text(text)
    text = unicode_safe(text, umode)

    if not text:
        log("WARNING: No text could be extracted from the document.", args.quiet)

    if args.output:
        out_path = Path(args.output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        # Always write files in UTF-8 regardless of terminal encoding
        out_path.write_text(text, encoding="utf-8")
        log(f"Saved to: {out_path}", args.quiet)
    else:
        try:
            print(text)
        except UnicodeEncodeError:
            # Last-resort fallback: re-run through replace mode and print
            log(
                "WARNING: Terminal cannot display some Unicode characters. "
                "Re-printing with --replace transliteration. "
                "Use -o output.txt to save the full Unicode text to a file.",
                quiet=False,
            )
            print(unicode_safe(text, "replace"))


# Importable API
__all__ = ["parse_document", "clean_text"]

if __name__ == "__main__":
    main()